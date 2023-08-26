# coding:utf-8

import pptx  # pip install python-pptx

p = pptx.Presentation()  # 生成ppt对象

layout = p.slide_layouts[8]  # 选择布局

# 0 title
# 1 title content
# 7
slide = p.slides.add_slide(layout)

p.save('test1.ppt')
