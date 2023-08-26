# coding:utf-8

import pptx

p = pptx.Presentation('test2.ppt')
for slide in p.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(shape.text_frame.text)
        if shape.has_table:
            for cell in shape.table.iter_cells():
                print(cell.text)

