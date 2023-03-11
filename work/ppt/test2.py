# coding:utf-8

import pptx
from pptx.util import Pt, Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

p = pptx.Presentation()
layout = p.slide_layouts[1]  # title content
slide = p.slides.add_slide(layout)

placeholder = slide.placeholders[1] # 0 title 1 content
#placeholder.text = '欢迎学习ppt制作\n欢迎学习python'

title = slide.placeholders[0]
title.text = '题目'

paragraph1 = placeholder.text_frame.add_paragraph()
paragraph1.text = '欢迎学习ppt制作'
paragraph1.bold = True
paragraph1.font.italic = True
paragraph1.font.size = Pt(16)
paragraph1.font.underline = True
paragraph1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

paragraph2 = placeholder.text_frame.add_paragraph()
paragraph2.text = '欢迎学习python'
paragraph2.font.size = Pt(32)
paragraph2.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT


layout = p.slide_layouts[6]  # title
slide = p.slides.add_slide(layout)
left = top = width = height = Inches(5)
box = slide.shapes.add_textbox(left, top, width, height)
para = box.text_frame.add_paragraph()

para.text = 'this is a para test'
para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
para.font.size = Pt(32)
para.font.color.rgb = RGBColor(255, 255, 0)
para.font.name = '微软雅黑'

layout = p.slide_layouts[1]
slide = p.slides.add_slide(layout)

rows = 10
cols = 2

left = top = Inches(2)
width = Inches(6.0)
height = Inches(1.0)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

for index, _ in enumerate(range(rows)):
    for sub_index in range(cols):
        table.cell(index, sub_index).text = '%s:%s' % (index, sub_index)

layout = p.slide_layouts[6]
slide = p.slides.add_slide(layout)

image = slide.shapes.add_picture(
    image_file='logo2020.png',
    left=Inches(1),
    top=Inches(1),
    width=Inches(6),
    height=Inches(4)
)

p.save('test2.ppt')
